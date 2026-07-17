#!/usr/bin/env python3
"""Generate leadership presentation for the UES Steam Safety Management Program."""

from io import BytesIO
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
ASSETS = ROOT / "docs" / "presentation" / "assets"
SCREENSHOTS = ROOT / "docs" / "presentation" / "screenshots"
OUTPUT = ROOT / "docs" / "presentation" / "UES-Steam-Safety-Program.pptx"

MAROON = RGBColor(0x50, 0x00, 0x00)
DARK = RGBColor(0x1E, 0x29, 0x3B)
GRAY = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0xB8, 0x3A, 0x3A)
LIGHT = RGBColor(0xFC, 0xE7, 0xE7)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
CREAM = RGBColor(0xFD, 0xF8, 0xF8)
SKY = RGBColor(0xE0, 0xF2, 0xFE)
GREEN_BG = RGBColor(0xEC, 0xFD, 0xF5)
AMBER_BG = RGBColor(0xFF, 0xFB, 0xEB)
TABLE_HEADER = MAROON
TABLE_ALT = RGBColor(0xF8, 0xFA, 0xFC)

FIELD_CHECKLIST_ROWS = [
    ("Valve body and bonnet", "No active corrosion; minor surface rust OK if no pitting"),
    ("Discharge pipe / stack", "Clear and unobstructed; properly supported; no standing water; drip pocket if applicable"),
    ("Drain / weep hole (if equipped)", "Open; no plugging or paint blockage"),
    ("Body and bonnet bolting", "All fasteners present and intact; no visible cracks or corrosion"),
    ("Nameplate legibility", "CDTP, serial, orifice, and manufacture year legible — record in work order"),
    ("Inlet / nozzle area", "No visible leakage or weeping at seat; dry mating surfaces"),
    ("Test gag / lifting lever", "Free of obstructions; lever not locked or wired shut (if equipped)"),
    ("Inlet isolation (if applicable)", "Valve in-line and not inadvertently isolated"),
]


class SlideCounter:
    def __init__(self) -> None:
        self.n = 0

    def mark(self, slide) -> None:
        self.n += 1
        box = slide.shapes.add_textbox(Inches(12.55), Inches(7.18), Inches(0.6), Inches(0.25))
        p = box.text_frame.paragraphs[0]
        p.text = str(self.n)
        style_paragraph(p, size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def set_slide_bg(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def style_paragraph(paragraph, *, size=14, bold=False, color=DARK, align=None) -> None:
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    if align is not None:
        paragraph.alignment = align


def add_header(slide, title: str, subtitle: str = "", *, compact: bool = False) -> float:
    title_size = 22 if compact else 28
    header = slide.shapes.add_textbox(Inches(0.65), Inches(0.32), Inches(12.1), Inches(0.65))
    hp = header.text_frame.paragraphs[0]
    hp.text = title
    style_paragraph(hp, size=title_size, bold=True, color=MAROON)

    line_y = 0.92 if compact else 1.02
    if subtitle:
        sp = slide.shapes.add_textbox(Inches(0.65), Inches(0.82 if compact else 0.9), Inches(12.1), Inches(0.38))
        s = sp.text_frame.paragraphs[0]
        s.text = subtitle
        style_paragraph(s, size=12 if compact else 13, color=GRAY)
        line_y = 1.15 if compact else 1.28

    if not compact:
        line = slide.shapes.add_shape(1, Inches(0.65), Inches(line_y), Inches(12.0), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT
        line.line.fill.background()
    return line_y + (0.08 if compact else 0.12)


def add_caveat_badge(slide, text: str, x: float = 9.8, y: float = 0.34) -> None:
    badge = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(2.9), Inches(0.38))
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(0xFE, 0xF3, 0xC7)
    badge.line.color.rgb = RGBColor(0xD9, 0x77, 0x06)
    tf = badge.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    style_paragraph(p, size=10, bold=True, color=RGBColor(0x92, 0x40, 0x0E), align=PP_ALIGN.CENTER)


def screenshot_stream(name: str) -> BytesIO | None:
    path = SCREENSHOTS / name
    if not path.exists():
        return None
    img = Image.open(path).convert("RGB")
    buf = BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, MAROON)
    accent = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.18))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    logo = ASSETS / "logo.png"
    if logo.exists():
        slide.shapes.add_picture(str(logo), Inches(6.15), Inches(0.55), width=Inches(1.05))

    title = slide.shapes.add_textbox(Inches(0.7), Inches(1.75), Inches(11.9), Inches(1.4))
    style_paragraph(title.text_frame.paragraphs[0], size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    title.text_frame.paragraphs[0].text = "Steam Safety Management Program"

    sub = slide.shapes.add_textbox(Inches(0.7), Inches(3.05), Inches(11.9), Inches(0.9))
    style_paragraph(sub.text_frame.paragraphs[0], size=20, color=LIGHT, align=PP_ALIGN.CENTER)
    sub.text_frame.paragraphs[0].text = "Pressure Safety Valve Compliance · Data · Maintenance Strategy"

    org = slide.shapes.add_textbox(Inches(0.7), Inches(3.95), Inches(11.9), Inches(0.6))
    style_paragraph(org.text_frame.paragraphs[0], size=18, color=RGBColor(0xE8, 0xC4, 0xC4), align=PP_ALIGN.CENTER)
    org.text_frame.paragraphs[0].text = "Texas A&M University · Utilities & Energy Services"

    chips = ["Field verified", "Excel master file", "Live dashboard", "Maintenance strategy"]
    chip_w = 2.55
    start_x = (13.333 - (chip_w * 4 + 0.2 * 3)) / 2
    for i, chip in enumerate(chips):
        x = start_x + i * (chip_w + 0.2)
        box = slide.shapes.add_shape(1, Inches(x), Inches(5.15), Inches(chip_w), Inches(0.48))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0x6B, 0x1A, 0x1A)
        box.line.color.rgb = ACCENT
        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = chip
        style_paragraph(p, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def add_section_slide(prs: Presentation, counter: SlideCounter, title: str, subtitle: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, MAROON)
    stripe = slide.shapes.add_shape(1, Inches(0), Inches(3.35), Inches(13.333), Inches(0.08))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ACCENT
    stripe.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.9), Inches(2.55), Inches(11.5), Inches(1.2))
    style_paragraph(box.text_frame.paragraphs[0], size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    box.text_frame.paragraphs[0].text = title

    if subtitle:
        sbox = slide.shapes.add_textbox(Inches(0.9), Inches(3.75), Inches(11.5), Inches(1))
        style_paragraph(sbox.text_frame.paragraphs[0], size=18, color=LIGHT, align=PP_ALIGN.CENTER)
        sbox.text_frame.paragraphs[0].text = subtitle
    counter.mark(slide)


def add_agenda_slide(prs: Presentation, counter: SlideCounter) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    top = add_header(slide, "Today's Briefing", "19 slides — context, program build, dashboard, strategy, next steps")

    sections = [
        ("I", "Context & Why It Matters", "Risk, regulatory foundation"),
        ("II", "Building the Program", "Field → Excel → validation path"),
        ("III", "Live Dashboard", "Production screenshots & capabilities"),
        ("IV", "Maintenance Strategy", "PM intervals, testing, checklist, tagging"),
        ("V", "Outcomes & Next Steps", "Value delivered, leadership decisions"),
    ]

    for i, (num, title, desc) in enumerate(sections):
        y = top + 0.2 + i * 1.05
        badge = slide.shapes.add_shape(1, Inches(0.75), Inches(y), Inches(0.55), Inches(0.55))
        badge.fill.solid()
        badge.fill.fore_color.rgb = MAROON
        badge.line.fill.background()
        bp = badge.text_frame.paragraphs[0]
        bp.text = num
        style_paragraph(bp, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        tb = slide.shapes.add_textbox(Inches(1.5), Inches(y - 0.02), Inches(10.5), Inches(0.75))
        tf = tb.text_frame
        p1 = tf.paragraphs[0]
        p1.text = title
        style_paragraph(p1, size=18, bold=True, color=DARK)
        p2 = tf.add_paragraph()
        p2.text = desc
        style_paragraph(p2, size=13, color=GRAY)

    counter.mark(slide)


def add_stat_cards(slide, stats: list[tuple[str, str, str]], top: float = 1.45) -> None:
    width = 2.85
    gap = 0.28
    start_x = (13.333 - (width * len(stats) + gap * (len(stats) - 1))) / 2
    for i, (value, label, hint) in enumerate(stats):
        x = start_x + i * (width + gap)
        card = slide.shapes.add_shape(1, Inches(x), Inches(top), Inches(width), Inches(1.35))
        card.fill.solid()
        card.fill.fore_color.rgb = CREAM
        card.line.color.rgb = ACCENT

        vbox = slide.shapes.add_textbox(Inches(x), Inches(top + 0.18), Inches(width), Inches(0.55))
        style_paragraph(vbox.text_frame.paragraphs[0], size=30, bold=True, color=MAROON, align=PP_ALIGN.CENTER)
        vbox.text_frame.paragraphs[0].text = value

        lbox = slide.shapes.add_textbox(Inches(x), Inches(top + 0.72), Inches(width), Inches(0.3))
        style_paragraph(lbox.text_frame.paragraphs[0], size=12, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        lbox.text_frame.paragraphs[0].text = label

        hbox = slide.shapes.add_textbox(Inches(x + 0.1), Inches(top + 1.0), Inches(width - 0.2), Inches(0.28))
        style_paragraph(hbox.text_frame.paragraphs[0], size=10, color=GRAY, align=PP_ALIGN.CENTER)
        hbox.text_frame.paragraphs[0].text = hint


def add_risk_cards_slide(prs: Presentation, counter: SlideCounter) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    top = add_header(slide, "Why Leadership Should Care", "Four risks this program directly addresses")

    cards = [
        ("⚠", "Regulatory", "Missed 3-year recert → TDLR / NBIC exposure"),
        ("👁", "Visibility", "Fleet status & spares were hard to see"),
        ("💰", "Cost", "Last-minute testing & unplanned outages"),
        ("📋", "Audit", "Scattered records, no defendable program"),
    ]
    positions = [(0.65, 0), (6.75, 0), (0.65, 2.85), (6.75, 2.85)]
    for (x, y_off), (icon, title, desc) in zip(positions, cards):
        y = top + y_off
        panel = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(5.95), Inches(2.55))
        panel.fill.solid()
        panel.fill.fore_color.rgb = CREAM if y_off == 0 and x < 5 else SKY if y_off == 0 else GREEN_BG
        panel.line.color.rgb = ACCENT

        ib = slide.shapes.add_textbox(Inches(x + 0.25), Inches(y + 0.2), Inches(0.5), Inches(0.45))
        style_paragraph(ib.text_frame.paragraphs[0], size=22, color=MAROON)
        ib.text_frame.paragraphs[0].text = icon

        tb = slide.shapes.add_textbox(Inches(x + 0.75), Inches(y + 0.2), Inches(5.0), Inches(1.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = title
        style_paragraph(p1, size=17, bold=True, color=DARK)
        p2 = tf.add_paragraph()
        p2.text = desc
        style_paragraph(p2, size=14, color=GRAY)

    counter.mark(slide)


def add_icon_row_slide(
    prs: Presentation,
    counter: SlideCounter,
    title: str,
    items: list[tuple[str, str, str]],
    subtitle: str = "",
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    top = add_header(slide, title, subtitle)

    n = len(items)
    width = min(2.9, 12.0 / n - 0.2)
    gap = 0.22
    start_x = (13.333 - (width * n + gap * (n - 1))) / 2
    for i, (icon, label, desc) in enumerate(items):
        x = start_x + i * (width + gap)
        card = slide.shapes.add_shape(1, Inches(x), Inches(top + 0.15), Inches(width), Inches(4.9))
        card.fill.solid()
        card.fill.fore_color.rgb = CREAM if i % 2 == 0 else SKY
        card.line.color.rgb = ACCENT

        ib = slide.shapes.add_textbox(Inches(x), Inches(top + 0.45), Inches(width), Inches(0.55))
        style_paragraph(ib.text_frame.paragraphs[0], size=28, color=MAROON, align=PP_ALIGN.CENTER)
        ib.text_frame.paragraphs[0].text = icon

        lb = slide.shapes.add_textbox(Inches(x + 0.1), Inches(top + 1.15), Inches(width - 0.2), Inches(0.7))
        style_paragraph(lb.text_frame.paragraphs[0], size=14, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        lb.text_frame.paragraphs[0].text = label

        db = slide.shapes.add_textbox(Inches(x + 0.12), Inches(top + 1.95), Inches(width - 0.24), Inches(2.8))
        tf = db.text_frame
        tf.word_wrap = True
        style_paragraph(tf.paragraphs[0], size=12, color=GRAY, align=PP_ALIGN.CENTER)
        tf.paragraphs[0].text = desc

        if i < n - 1:
            ar = slide.shapes.add_textbox(Inches(x + width + 0.02), Inches(top + 2.2), Inches(0.2), Inches(0.3))
            style_paragraph(ar.text_frame.paragraphs[0], size=18, color=ACCENT, align=PP_ALIGN.CENTER)
            ar.text_frame.paragraphs[0].text = "→"

    counter.mark(slide)


def style_table_cell(cell, text: str, *, header: bool = False, alt: bool = False, size: int = 12) -> None:
    cell.text = text
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if header:
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER
    elif alt:
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_ALT
    else:
        cell.fill.background()
    for paragraph in cell.text_frame.paragraphs:
        style_paragraph(paragraph, size=size, bold=header, color=WHITE if header else DARK)


def add_table_slide(
    prs: Presentation,
    counter: SlideCounter,
    title: str,
    headers: list[str],
    rows: list[list[str]],
    subtitle: str = "",
    col_widths: list[float] | None = None,
    caveat_badge: str = "",
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    top = add_header(slide, title, subtitle)
    if caveat_badge:
        add_caveat_badge(slide, caveat_badge)

    table_shape = slide.shapes.add_table(
        len(rows) + 1, len(headers), Inches(0.65), Inches(top), Inches(12.0), Inches(5.6)
    )
    table = table_shape.table
    if col_widths:
        for idx, width in enumerate(col_widths):
            table.columns[idx].width = Inches(width)
    for col, header in enumerate(headers):
        style_table_cell(table.cell(0, col), header, header=True, size=12)
    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            style_table_cell(table.cell(r, c), value, alt=r % 2 == 0, size=11)
    counter.mark(slide)


def add_bullet_slide(
    prs: Presentation,
    counter: SlideCounter,
    title: str,
    bullets: list[str],
    subtitle: str = "",
    font_size: int = 17,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    top = add_header(slide, title, subtitle)
    for i, bullet in enumerate(bullets):
        y = top + 0.15 + i * 0.95
        marker = slide.shapes.add_shape(1, Inches(0.75), Inches(y + 0.1), Inches(0.14), Inches(0.14))
        marker.fill.solid()
        marker.fill.fore_color.rgb = ACCENT
        marker.line.fill.background()
        box = slide.shapes.add_textbox(Inches(1.05), Inches(y), Inches(11.5), Inches(0.8))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = bullet
        style_paragraph(p, size=font_size, color=DARK)
    counter.mark(slide)


def add_two_column_cards(
    prs: Presentation,
    counter: SlideCounter,
    title: str,
    left_title: str,
    left: list[str],
    right_title: str,
    right: list[str],
    subtitle: str = "",
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    top = add_header(slide, title, subtitle)
    for col_x, col_title, items, bg in ((0.65, left_title, left, CREAM), (6.75, right_title, right, SKY)):
        panel = slide.shapes.add_shape(1, Inches(col_x), Inches(top), Inches(5.95), Inches(5.55))
        panel.fill.solid()
        panel.fill.fore_color.rgb = bg
        panel.line.color.rgb = ACCENT
        label = slide.shapes.add_textbox(Inches(col_x + 0.2), Inches(top + 0.15), Inches(5.5), Inches(0.35))
        style_paragraph(label.text_frame.paragraphs[0], size=15, bold=True, color=MAROON)
        label.text_frame.paragraphs[0].text = col_title
        box = slide.shapes.add_textbox(Inches(col_x + 0.25), Inches(top + 0.55), Inches(5.45), Inches(4.8))
        tf = box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {item}"
            style_paragraph(p, size=14, color=DARK)
            p.space_after = Pt(8)
    counter.mark(slide)


def add_journey_slide(prs: Presentation, counter: SlideCounter) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    top = add_header(slide, "Program Development Path", "From field discovery to a sustainable maintenance program")
    phases = [
        ("1", "Field Data\nCollection", "Site walks\nNameplate photos\nValve verification"),
        ("2", "Master Excel\nRegister", "Sort & normalize\nTrack due dates\nInventory IDs"),
        ("3", "Digital\nDashboard", "Fleet KPIs\nDrill-down views\nExcel exports"),
        ("4", "Maintenance\nStrategy", "PM intervals\nSpare pooling\nCode alignment"),
    ]
    start_x, width, gap = 0.55, 2.95, 0.35
    for i, (num, phase_title, desc) in enumerate(phases):
        x = start_x + i * (width + gap)
        card = slide.shapes.add_shape(1, Inches(x), Inches(top + 0.15), Inches(width), Inches(4.95))
        card.fill.solid()
        card.fill.fore_color.rgb = CREAM
        card.line.color.rgb = ACCENT
        circle = slide.shapes.add_shape(1, Inches(x + 0.2), Inches(top + 0.35), Inches(0.55), Inches(0.55))
        circle.fill.solid()
        circle.fill.fore_color.rgb = MAROON
        circle.line.fill.background()
        np = circle.text_frame.paragraphs[0]
        np.text = num
        style_paragraph(np, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        title_box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(top + 1.05), Inches(width - 0.3), Inches(1.0))
        style_paragraph(title_box.text_frame.paragraphs[0], size=16, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        title_box.text_frame.paragraphs[0].text = phase_title
        desc_box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(top + 2.2), Inches(width - 0.3), Inches(2.5))
        style_paragraph(desc_box.text_frame.paragraphs[0], size=13, color=GRAY, align=PP_ALIGN.CENTER)
        desc_box.text_frame.paragraphs[0].text = desc
        if i < len(phases) - 1:
            arrow = slide.shapes.add_textbox(Inches(x + width + 0.05), Inches(top + 2.35), Inches(0.3), Inches(0.4))
            style_paragraph(arrow.text_frame.paragraphs[0], size=24, color=ACCENT, align=PP_ALIGN.CENTER)
            arrow.text_frame.paragraphs[0].text = "→"
    counter.mark(slide)


def add_screenshot_slide(
    prs: Presentation,
    counter: SlideCounter,
    title: str,
    image_name: str,
    caption: str,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    # Compact title — image dominates the slide
    tb = slide.shapes.add_textbox(Inches(0.45), Inches(0.18), Inches(10.5), Inches(0.42))
    style_paragraph(tb.text_frame.paragraphs[0], size=20, bold=True, color=MAROON)
    tb.text_frame.paragraphs[0].text = title

    stream = screenshot_stream(image_name)
    if stream:
        slide.shapes.add_picture(stream, Inches(0.2), Inches(0.55), width=Inches(12.9))

    cap = slide.shapes.add_shape(1, Inches(0.2), Inches(6.95), Inches(12.9), Inches(0.38))
    cap.fill.solid()
    cap.fill.fore_color.rgb = RGBColor(0xF1, 0xF5, 0xF9)
    cap.line.fill.background()
    cp = cap.text_frame.paragraphs[0]
    cp.text = caption
    style_paragraph(cp, size=10, color=GRAY, align=PP_ALIGN.CENTER)
    cap.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    counter.mark(slide)


def add_field_checklist_slide(prs: Presentation, counter: SlideCounter) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    top = add_header(
        slide,
        "Pre-Test Visual Inspection — Field Checklist",
        "Section 4.1 · Complete before every removal or in-place test · File in CMMS work order",
    )
    headers = ["Checkpoint", "Accept / Reject criteria", "Pass", "Fail", "N/A"]
    rows = [[cp, crit, "☐", "☐", "☐"] for cp, crit in FIELD_CHECKLIST_ROWS]
    table_shape = slide.shapes.add_table(
        len(rows) + 1, len(headers), Inches(0.65), Inches(top + 0.05), Inches(12.0), Inches(5.35)
    )
    table = table_shape.table
    for idx, width in enumerate([2.6, 7.0, 0.75, 0.75, 0.75]):
        table.columns[idx].width = Inches(width)
    for col, header in enumerate(headers):
        style_table_cell(table.cell(0, col), header, header=True, size=11)
    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            style_table_cell(cell, value, alt=r % 2 == 0, size=9 if c == 1 else 10)
            cell.text_frame.word_wrap = True
            if c >= 2:
                for p in cell.text_frame.paragraphs:
                    p.alignment = PP_ALIGN.CENTER
    note = slide.shapes.add_shape(1, Inches(0.65), Inches(6.62), Inches(12.0), Inches(0.42))
    note.fill.solid()
    note.fill.fore_color.rgb = AMBER_BG
    note.line.color.rgb = ACCENT
    tf = note.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "Technician: _______________   Date: __________   WO #: __________   Equipment / Location: _______________"
    style_paragraph(p, size=10, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    counter.mark(slide)


def add_tag_diagram_slide(prs: Presentation, counter: SlideCounter) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    top = add_header(slide, "Valve Identification & Labeling", "Section 8 · Read the tag left to right — each segment has a fixed meaning")
    diagram = ASSETS / "tag-diagram.png"
    if diagram.exists():
        slide.shapes.add_picture(str(diagram), Inches(0.65), Inches(top), width=Inches(12.0))
    note = slide.shapes.add_shape(1, Inches(0.65), Inches(6.45), Inches(12.0), Inches(0.55))
    note.fill.solid()
    note.fill.fore_color.rgb = AMBER_BG
    note.line.color.rgb = ACCENT
    tf = note.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "Physical tag is two-sided: Position Tag + Inventory ID + Serial Number + Service + Location"
    style_paragraph(p, size=13, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    counter.mark(slide)


def add_decision_columns_slide(prs: Presentation, counter: SlideCounter) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    top = add_header(slide, "Overhaul vs. Repair vs. Replace", "Decision guide for Setpoint and procurement")

    columns = [
        ("🔧", "Overhaul", "Scheduled 3-year send-out", "Full disassembly, lapping, spring check, certified pop test"),
        ("🛠", "Repair", "Limited defect found", "VR-stamp scope + re-test before return to service"),
        ("🔄", "Replace", "Body damage or obsolete", "New valve when economics or risk favor replacement"),
    ]
    col_w = 3.85
    start_x = 0.65
    for i, (icon, title, when, detail) in enumerate(columns):
        x = start_x + i * (col_w + 0.25)
        panel = slide.shapes.add_shape(1, Inches(x), Inches(top + 0.1), Inches(col_w), Inches(5.2))
        panel.fill.solid()
        panel.fill.fore_color.rgb = [CREAM, SKY, GREEN_BG][i]
        panel.line.color.rgb = ACCENT
        ib = slide.shapes.add_textbox(Inches(x), Inches(top + 0.35), Inches(col_w), Inches(0.5))
        style_paragraph(ib.text_frame.paragraphs[0], size=26, color=MAROON, align=PP_ALIGN.CENTER)
        ib.text_frame.paragraphs[0].text = icon
        tb = slide.shapes.add_textbox(Inches(x + 0.15), Inches(top + 0.95), Inches(col_w - 0.3), Inches(4.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = title
        style_paragraph(p1, size=17, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        p2 = tf.add_paragraph()
        p2.text = when
        style_paragraph(p2, size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        p3 = tf.add_paragraph()
        p3.text = detail
        style_paragraph(p3, size=13, color=GRAY, align=PP_ALIGN.CENTER)

    foot = slide.shapes.add_textbox(Inches(0.65), Inches(6.55), Inches(12), Inches(0.35))
    style_paragraph(foot.text_frame.paragraphs[0], size=11, color=GRAY, align=PP_ALIGN.CENTER)
    foot.text_frame.paragraphs[0].text = "Spares in storage >5 years: re-test before install per API RP 576"
    counter.mark(slide)


def add_capability_chips(slide, top: float, items: list[tuple[str, str]]) -> float:
    n = len(items)
    chip_w = min(2.95, 12.0 / n - 0.12)
    gap = 0.12
    start_x = (13.333 - (chip_w * n + gap * (n - 1))) / 2
    for i, (label, desc) in enumerate(items):
        x = start_x + i * (chip_w + gap)
        box = slide.shapes.add_shape(1, Inches(x), Inches(top), Inches(chip_w), Inches(0.72))
        box.fill.solid()
        box.fill.fore_color.rgb = CREAM if i % 2 == 0 else SKY
        box.line.color.rgb = ACCENT
        tf = box.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = label
        style_paragraph(p1, size=11, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        p2 = tf.add_paragraph()
        p2.text = desc
        style_paragraph(p2, size=9, color=GRAY, align=PP_ALIGN.CENTER)
    return top + 0.82


def add_dashboard_hero_slide(prs: Presentation, counter: SlideCounter) -> None:
    """Dashboard capabilities as a slim chip row above a full-bleed screenshot."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    top = add_header(
        slide,
        "Phase 3 — Live Compliance Dashboard",
        "reliability-and-compliance-dashboar.vercel.app",
        compact=True,
    )
    img_top = add_capability_chips(slide, top + 0.05, [
        ("Fleet KPIs", "Total, OOS, Overdue, Compliant %"),
        ("Drill-down", "KPI → valve list with Inventory ID"),
        ("Hierarchy", "Equipment → Location → PSV"),
        ("Export", "5-sheet Excel for meetings"),
    ])

    stream = screenshot_stream("01-dashboard.png")
    if stream:
        slide.shapes.add_picture(stream, Inches(0.2), Inches(img_top), width=Inches(12.9))

    cap = slide.shapes.add_shape(1, Inches(0.2), Inches(6.88), Inches(12.9), Inches(0.38))
    cap.fill.solid()
    cap.fill.fore_color.rgb = RGBColor(0xF1, 0xF5, 0xF9)
    cap.line.fill.background()
    cp = cap.text_frame.paragraphs[0]
    cp.text = "Live production — 51 PSVs, 78% compliant, fleet-wide equipment view"
    style_paragraph(cp, size=10, color=GRAY, align=PP_ALIGN.CENTER)
    cap.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    counter.mark(slide)


def add_outcomes_and_next_steps_slide(prs: Presentation, counter: SlideCounter) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    top = add_header(slide, "Program Outcomes & Next Steps", "Value delivered · leadership decisions to sustain the program")

    for col_x, col_title, items, bg in (
        (0.65, "Operational value", [
            "Complete fleet inventory & visible compliance by site",
            "Traceable audit history & faster Setpoint planning",
        ], CREAM),
        (6.75, "Strategic value", [
            "Code-aligned strategy doc (ASME, NBIC, API RP 576)",
            "CMMS foundation, spare-pool path, sustainable UES ownership",
        ], SKY),
    ):
        panel = slide.shapes.add_shape(1, Inches(col_x), Inches(top), Inches(5.95), Inches(2.35))
        panel.fill.solid()
        panel.fill.fore_color.rgb = bg
        panel.line.color.rgb = ACCENT
        label = slide.shapes.add_textbox(Inches(col_x + 0.2), Inches(top + 0.12), Inches(5.5), Inches(0.3))
        style_paragraph(label.text_frame.paragraphs[0], size=14, bold=True, color=MAROON)
        label.text_frame.paragraphs[0].text = col_title
        box = slide.shapes.add_textbox(Inches(col_x + 0.25), Inches(top + 0.48), Inches(5.45), Inches(1.65))
        tf = box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {item}"
            style_paragraph(p, size=13, color=DARK)
            p.space_after = Pt(8)

    steps_top = top + 2.65
    steps_label = slide.shapes.add_textbox(Inches(0.65), Inches(steps_top), Inches(12), Inches(0.3))
    style_paragraph(steps_label.text_frame.paragraphs[0], size=14, bold=True, color=MAROON)
    steps_label.text_frame.paragraphs[0].text = "Recommended next steps"

    for i, bullet in enumerate([
        "Endorse Maintenance Strategy and PM frequencies",
        "Approve spare-acquisition priority list",
        "Pilot test-and-reuse on 125 psi valves (3–5 unit pilot)",
        "Fleet-wide physical tagging — Position Tag + Inventory ID",
    ]):
        y = steps_top + 0.42 + i * 0.72
        marker = slide.shapes.add_shape(1, Inches(0.75), Inches(y + 0.08), Inches(0.14), Inches(0.14))
        marker.fill.solid()
        marker.fill.fore_color.rgb = ACCENT
        marker.line.fill.background()
        box = slide.shapes.add_textbox(Inches(1.05), Inches(y), Inches(11.5), Inches(0.65))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = bullet
        style_paragraph(p, size=15, color=DARK)

    counter.mark(slide)


def add_replacement_cycle_slide(prs: Presentation, counter: SlideCounter) -> None:
    """3-year location-based replacement rotation to spread capital costs."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    top = add_header(
        slide,
        "3-Year Replacement Cycle by Location",
        "Stagger replacements across sites so capital costs are predictable — not concentrated in one year",
    )

    why = slide.shapes.add_textbox(Inches(0.65), Inches(top), Inches(12.0), Inches(1.05))
    tf = why.text_frame
    tf.word_wrap = True
    for i, line in enumerate([
        "Without a rotation plan, valves installed together age out together — creating a single-year budget spike.",
        "Assign each building or boiler location to a Year 1, 2, or 3 bucket so replacement and recert spend is level-loaded.",
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        style_paragraph(p, size=13, color=DARK)
        p.space_after = Pt(4)

    years = [
        ("Year 1", "Rotation group A", [
            "Assign ~⅓ of locations (by valve count)",
            "Commercial 150 psi replacements due this cycle",
            "Setpoint send-outs aligned to spare swaps",
        ], CREAM),
        ("Year 2", "Rotation group B", [
            "Next third of campus / plant locations",
            "Balance workload with Year 1 completions",
            "Update dashboard due dates after each swap",
        ], SKY),
        ("Year 3", "Rotation group C", [
            "Remaining locations complete the cycle",
            "Cycle restarts — locations return to Year 1",
            "One due-date logic across the whole fleet",
        ], GREEN_BG),
    ]
    col_w = 3.85
    card_top = top + 1.2
    for i, (year, group, items, bg) in enumerate(years):
        x = 0.65 + i * (col_w + 0.25)
        panel = slide.shapes.add_shape(1, Inches(x), Inches(card_top), Inches(col_w), Inches(3.55))
        panel.fill.solid()
        panel.fill.fore_color.rgb = bg
        panel.line.color.rgb = ACCENT

        yb = slide.shapes.add_shape(1, Inches(x + 0.15), Inches(card_top + 0.15), Inches(col_w - 0.3), Inches(0.42))
        yb.fill.solid()
        yb.fill.fore_color.rgb = MAROON
        yb.line.fill.background()
        yp = yb.text_frame.paragraphs[0]
        yp.text = year
        style_paragraph(yp, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        gt = slide.shapes.add_textbox(Inches(x + 0.15), Inches(card_top + 0.65), Inches(col_w - 0.3), Inches(0.35))
        style_paragraph(gt.text_frame.paragraphs[0], size=12, bold=True, color=MAROON, align=PP_ALIGN.CENTER)
        gt.text_frame.paragraphs[0].text = group

        box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(card_top + 1.05), Inches(col_w - 0.4), Inches(2.3))
        btf = box.text_frame
        btf.word_wrap = True
        for j, item in enumerate(items):
            p = btf.paragraphs[0] if j == 0 else btf.add_paragraph()
            p.text = f"• {item}"
            style_paragraph(p, size=11, color=DARK)
            p.space_after = Pt(5)

    foot = slide.shapes.add_shape(1, Inches(0.65), Inches(6.35), Inches(12.0), Inches(0.75))
    foot.fill.solid()
    foot.fill.fore_color.rgb = AMBER_BG
    foot.line.color.rgb = ACCENT
    ftf = foot.text_frame
    ftf.word_wrap = True
    ftf.vertical_anchor = MSO_ANCHOR.MIDDLE
    fp = ftf.paragraphs[0]
    fp.text = (
        "How to build the plan: use dashboard valve counts and due dates per location · "
        "prioritize failed annual lever tests and overdue recerts into the nearest rotation year · "
        "leadership approves a 3-year capital forecast (~⅓ of fleet cost per year)"
    )
    style_paragraph(fp, size=11, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    counter.mark(slide)


def add_for_against_slide(
    prs: Presentation,
    counter: SlideCounter,
    title: str,
    for_title: str,
    for_points: list[str],
    against_title: str,
    against_points: list[str],
    subtitle: str = "",
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    top = add_header(slide, title, subtitle)

    for col_x, col_title, items, bg in ((0.65, for_title, for_points, GREEN_BG), (6.75, against_title, against_points, AMBER_BG)):
        panel = slide.shapes.add_shape(1, Inches(col_x), Inches(top), Inches(5.95), Inches(5.55))
        panel.fill.solid()
        panel.fill.fore_color.rgb = bg
        panel.line.color.rgb = ACCENT
        label = slide.shapes.add_textbox(Inches(col_x + 0.2), Inches(top + 0.15), Inches(5.5), Inches(0.35))
        style_paragraph(label.text_frame.paragraphs[0], size=14, bold=True, color=MAROON)
        label.text_frame.paragraphs[0].text = col_title
        box = slide.shapes.add_textbox(Inches(col_x + 0.25), Inches(top + 0.55), Inches(5.45), Inches(4.8))
        tf = box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {item}"
            style_paragraph(p, size=12, color=DARK)
            p.space_after = Pt(6)
    counter.mark(slide)


def build() -> Path:
    tag_script = ROOT / "scripts" / "generate_tag_diagram.py"
    if tag_script.exists():
        import runpy
        runpy.run_path(str(tag_script), run_name="__not_main__")

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    counter = SlideCounter()

    add_title_slide(prs)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_header(slide, "Executive Summary", "End-to-end program built from the ground up")
    add_stat_cards(slide, [
        ("4 Phases", "Program path", "Field → Excel → Dashboard → Strategy"),
        ("51+ PSVs", "Tracked fleet", "Boilers, HRSG, HHWB, DHWB"),
        ("3 Years", "Recert cycle", "NBIC / ASME aligned"),
        ("1 Source", "System of record", "Dashboard + master Excel"),
    ])
    box = slide.shapes.add_textbox(Inches(0.85), Inches(3.05), Inches(11.8), Inches(3.2))
    tf = box.text_frame
    tf.word_wrap = True
    for i, point in enumerate([
        "Walked boiler rooms and steam systems — photographed nameplates and verified installed valves",
        "Built a comprehensive Excel master register, then a live compliance dashboard for the full fleet",
        "Authored a formal PSV/SRV Maintenance Strategy aligned to ASME, NBIC, API RP 576, and Texas code",
        "Outcome: visible compliance, proactive due-date management, and a defendable maintenance program",
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"✓  {point}"
        style_paragraph(p, size=16, color=DARK)
        p.space_after = Pt(12)
    counter.mark(slide)

    add_bullet_slide(
        prs, counter,
        "Why Leadership Should Care",
        [
            "Boiler PSVs are code-required — missed recertification creates regulatory risk",
            "Valve status, spare coverage, and due dates were hard to see fleet-wide",
            "Unplanned failures and rushed testing drive cost, outages, and audit risk",
            "A structured program cuts surprises and improves capital planning",
            "This reliability work protects people, equipment, and compliance standing",
        ],
        subtitle="Risk reduction through visibility and a sustainable maintenance program",
        font_size=16,
    )

    add_bullet_slide(
        prs, counter,
        "Regulatory & Standards Foundation",
        [
            "ASME Section I — boiler safety valve design, capacity, and manual-lift requirements",
            "NBIC Part 2 — in-service inspection and maximum 3-year test intervals",
            "API RP 576 — industry best practice for inspection, testing, and repair of relief devices",
            "16 TAC Chapter 65 & Texas Health & Safety Code Ch. 755 — state boiler program requirements",
            "Maintenance Strategy document ties every PM activity back to these references",
        ],
        subtitle="Program decisions are code-informed, not ad hoc",
        font_size=15,
    )

    add_section_slide(prs, counter, "Building the Program", "Field data → Master register → Dashboard → Maintenance strategy")
    add_journey_slide(prs, counter)

    add_bullet_slide(
        prs, counter,
        "Phase 1 — Field Data Collection",
        [
            "Walked every boiler room to locate all relief valves and protected points",
            "Captured nameplate photos: serial, set pressure (CDTP), orifice, NB number",
            "Verified what was actually installed vs. what drawings or records showed",
            "Documented spare availability, valve orientation, discharge path, and visible condition",
            "Flagged gaps: missing spares, illegible nameplates, weeping seats, isolation concerns",
        ],
        subtitle="Ground truth before any database or dashboard",
        font_size=15,
    )

    add_two_column_cards(
        prs, counter,
        "What We Captured in the Field",
        "Equipment & location context",
        [
            "Boiler / HRSG / commercial unit identity",
            "Protected location (steam drum, economizer, etc.)",
            "Position on the system",
            "Installed vs. spare status",
        ],
        "Valve identity & condition",
        [
            "Serial number & manufacturer",
            "Set pressure and capacity",
            "National Board number",
            "Nameplate photo evidence",
            "Visible leakage, corrosion, or damage",
        ],
    )

    add_bullet_slide(
        prs, counter,
        "Phase 2 — Master Excel Register",
        [
            "Built a comprehensive Excel master file as the first system of record",
            "Sorted and normalized field data across equipment, locations, and serial numbers",
            "Tracked install dates, service history, spare availability, and recert due dates",
            "Assigned Inventory IDs and position-oriented tags for procurement and CMMS use",
            "Used Excel to validate completeness before investing in a digital platform",
            "Master file still supports bulk reporting and leadership reviews",
        ],
        subtitle="The bridge between field discovery and digital management",
        font_size=14,
    )

    add_section_slide(prs, counter, "Phase 3 — Digital Compliance Dashboard", "Operational visibility for the whole fleet")

    add_bullet_slide(
        prs, counter,
        "Dashboard Capabilities",
        [
            "Site-wide KPIs: Total PSVs, Installed, Out for Service, Overdue, Compliant %",
            "Equipment → Location → PSV hierarchy matching how operators think about the plant",
            "Automatic 3-year due-date calculation from install or on-site service history",
            "Inventory ID visible on locations, KPI views, and exports",
            "Status history and separate repair/overhaul history per valve",
            "Excel report download: All PSVs, Installed, Out for Service, Overdue, Upcoming Due",
        ],
        subtitle="reliability-and-compliance-dashboar.vercel.app",
        font_size=14,
    )

    add_screenshot_slide(prs, counter, "Dashboard — Fleet Compliance at a Glance", "01-dashboard.png",
                         "Live production — 51+ PSVs tracked, fleet-wide equipment view")
    add_screenshot_slide(prs, counter, "KPI Drill-Down", "02-kpi-modal.png",
                         "Click any KPI to see the exact valves driving that metric")
    add_screenshot_slide(prs, counter, "Equipment & Location Views", "03-equipment.png",
                         "Scoped view per boiler with location-level Inventory ID and installed valve due dates")
    add_screenshot_slide(prs, counter, "Valve Record Detail", "05-psv-detail.png",
                         "Full datasheet, compliance clock, status history, and repair/overhaul log")

    add_section_slide(prs, counter, "Phase 4 — Maintenance Strategy", "Formal PSV/SRV program document for UES operations")

    add_table_slide(
        prs, counter,
        "PM Intervals & Triggers",
        ["Activity", "Frequency", "Code / reference", "Notes"],
        [
            ["Formal pop test & overhaul", "Every 3 years", "NBIC Part 2 §2.5.8; ASME Sec. I; 16 TAC Ch. 65", "Code minimum — Setpoint send-out or in-place"],
            ["Visual field walkdown", "Semi-annual", "API RP 576 §8 & §9", "Eyes-on only — recommended UES program addition"],
            ["Manual lever test", "Annual (optional)", "ASME Sec. I PG-73.1.3", "Management decision — see Section 3"],
            ["Condition-based inspection", "Immediately", "API RP 576 §9", "Leakage, chattering, corrosion, unplanned actuation"],
        ],
        subtitle="Maximum intervals shown — UES should complete before due date",
        col_widths=[2.4, 1.7, 3.2, 4.7],
    )

    add_table_slide(
        prs, counter,
        "Testing Methods — Send-Out vs. In-Place",
        ["Factor", "Send-out (Setpoint shop)", "In-place (Setpoint on-site)"],
        [
            ["Who performs", "Setpoint certified shop technicians", "Same Setpoint technicians, on-site"],
            ["Certification", "Full Setpoint test certificate", "Full certificate — equivalent"],
            ["Scope of work", "Bench overhaul: disassembly, lapping, pop test", "As-found lift; reseat & seat tightness check"],
            ["Internal assessment", "Complete — worn parts replaced", "Limited — functional test only"],
            ["Seat damage risk", "Lower — controlled bench environment", "Higher — especially older or corroded valves"],
            ["Best suited for", "Preferred when spare is available", "Interim until spare is acquired"],
        ],
        subtitle="Long-term goal: spare coverage for all valves to enable send-out testing fleet-wide",
        col_widths=[2.0, 5.0, 5.0],
    )

    add_field_checklist_slide(prs, counter)
    add_tag_diagram_slide(prs, counter)

    add_for_against_slide(
        prs, counter,
        "Manual Lever Testing — Management Decision Required",
        "Arguments for",
        [
            "PG-73.1.3 requires capability to manually lift at ≥75% of set pressure on Section I boilers",
            "Low-cost field check that disc is not stuck to seat from corrosion or scale",
            "Brief test at operating pressure confirms disc and seat are functional",
        ],
        "Arguments against / risks",
        [
            "Can damage seating surfaces — especially older valves or soft seats",
            "Post-test seat leakage may require immediate removal and bench overhaul",
            "API RP 576 does not recommend routine lever tests as substitute for formal pop testing",
            "Low incremental benefit if formal test is already scheduled",
        ],
        subtitle="Section 3 — both testing methods performed by Setpoint; certification equivalent",
    )

    add_bullet_slide(
        prs, counter,
        "Overhaul vs. Repair vs. Replace",
        [
            "Overhaul: full disassembly, lapping, and certified pop test — bundled with send-out",
            "Repair: limited-scope fix — must be VR-stamp work with re-test before return",
            "Replace when: body damage, obsolete design, or high-criticality risk favors new",
            "Spares in storage >5 years: re-test before installation per API RP 576 guidance",
            "Decision criteria documented so procurement and operations align on the same rules",
        ],
        subtitle="Decision guide for Setpoint and procurement",
        font_size=15,
    )

    add_bullet_slide(
        prs, counter,
        "Commercial Boiler Safety Valves (150 PSI)",
        [
            "Section 6 — ASME Section IV heating-boiler devices; different class from Section I steam PRVs",
            "F275464 (~$351/valve) replaces discontinued 0275464 — use ~$350–360 for planning",
            "Use-and-trash is correct: repair cost exceeds ~60% of replacement at this price point",
            "Annual lever test required per Watts ES-174A/740 — replace immediately on failed reseat",
            "Recommend 3–5 year replacement cycle; align with fleet 3-year PSV cadence (Section 6.4)",
        ],
        subtitle="Cost-driven maintenance logic — test-and-reuse is not worth pursuing for this valve class",
        font_size=14,
    )

    add_replacement_cycle_slide(prs, counter)

    add_bullet_slide(
        prs, counter,
        "Highest-Impact Program Improvements",
        [
            "Acquire spares for valves with none — enables full bench testing",
            "Rank valves by criticality to prioritize capital spend",
            "Adopt semi-annual visual walkdowns as standard UES practice between formal tests",
            "Keep dashboard and master register as the single source of truth",
            "Management decision needed on routine manual lever testing — risk vs. PG-73.1.3 compliance",
        ],
        subtitle="Section 5.5 — spare coverage is the single highest-impact improvement",
        font_size=15,
    )

    add_two_column_cards(
        prs, counter,
        "Program Outcomes & Value Delivered",
        "Operational value",
        [
            "Complete valve inventory across steam systems",
            "Visible compliance status by equipment and site",
            "Traceable history for audits and incidents",
            "Faster planning for Setpoint outages and swaps",
            "Reduced reliance on memory and scattered files",
        ],
        "Strategic value",
        [
            "Code-aligned maintenance strategy document",
            "Foundation for CMMS integration and spare pooling",
            "Data to support capital requests for spares",
            "Leadership-ready reporting via dashboard & Excel",
            "Sustainable program ownership within UES",
        ],
    )

    add_bullet_slide(
        prs, counter,
        "Recommended Next Steps",
        [
            "Leadership endorsement of the Maintenance Strategy and PM frequencies",
            "Decision on manual lever testing policy (Section 3 — management choice)",
            "Approve 3-year location rotation plan for commercial replacements and recert spend",
            "Approve spare-acquisition prioritization for no-spare locations",
            "Continue populating dashboard with remaining field discoveries (walkdowns in progress)",
            "Integrate Position Tags and Inventory IDs fleet-wide on physical valve tags",
        ],
        subtitle="Leadership decisions to sustain the program",
        font_size=15,
    )

    add_section_slide(prs, counter, "Discussion", "Steam Safety Management Program · Texas A&M UES · Reliability & Compliance")

    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    path = build()
    print(f"Created: {path}")
